import os
import json
import pandas as pd
import pdfplumber
from pptx import Presentation

class FileIngestion:
    def __init__(self, data_folder="data", verbose=False):
        """Initialize the ingestion class with the specified data folder."""
        self.data_folder = data_folder
        self.verbose = verbose  # Controls final print output

    def load_json_files(self):
        """Loads all JSON files in the data folder and returns a structured dictionary."""
        json_data = {}
        json_files = [f for f in os.listdir(self.data_folder) if f.lower().endswith(".json")]

        for file in json_files:
            file_path = os.path.join(self.data_folder, file)
            try:
                with open(file_path, "r", encoding="utf-8") as f:
                    data = json.load(f)
                    json_data[file.replace(".json", "")] = data
            except Exception as e:
                print(f"⚠️ Error loading JSON {file}: {e}")

        return json_data if json_data else {}

    def load_csv_files(self):
        """Loads all CSV files in the data folder and returns a dictionary of DataFrames."""
        csv_data = {}
        csv_files = [f for f in os.listdir(self.data_folder) if f.lower().endswith(".csv")]

        for file in csv_files:
            file_path = os.path.join(self.data_folder, file)
            try:
                df = pd.read_csv(file_path, encoding="utf-8-sig", low_memory=False)
                csv_data[file.replace(".csv", "")] = df
            except Exception as e:
                print(f"⚠️ Error loading CSV {file}: {e}")

        return csv_data if csv_data else {}

    def load_pdf_files(self):
        """Loads all PDF files and extracts tables & text data."""
        pdf_data = {}
        pdf_files = [f for f in os.listdir(self.data_folder) if f.lower().endswith(".pdf")]

        for file in pdf_files:
            file_path = os.path.join(self.data_folder, file)
            try:
                text_data, table_data = [], []
                with pdfplumber.open(file_path) as pdf:
                    for page in pdf.pages:
                        text = page.extract_text()
                        if text:
                            text_data.append(text)

                        tables = page.extract_tables()
                        for table in tables:
                            if table:
                                df = pd.DataFrame(table)
                                df.columns = df.iloc[0]  # Set first row as headers
                                df = df[1:].reset_index(drop=True)
                                table_data.append(df)

                pdf_data[file.replace(".pdf", "")] = {
                    "text": text_data or None,
                    "tables": table_data or None
                }
            except Exception as e:
                print(f"⚠️ Error loading PDF {file}: {e}")

        return pdf_data if pdf_data else {}

    def load_pptx_files(self):
        """Loads all PPTX files and extracts slide content and tables."""
        pptx_data = {}
        pptx_files = [f for f in os.listdir(self.data_folder) if f.lower().endswith(".pptx")]

        for file in pptx_files:
            file_path = os.path.join(self.data_folder, file)
            try:
                prs = Presentation(file_path)
                slide_data = []

                for slide_num, slide in enumerate(prs.slides, start=1):
                    slide_text = "\n".join([shape.text for shape in slide.shapes if hasattr(shape, "text")])
                    table_data = []

                    for shape in slide.shapes:
                        if hasattr(shape, "table"):
                            df = pd.DataFrame([[cell.text for cell in row.cells] for row in shape.table.rows])
                            if not df.empty:
                                df.columns = df.iloc[0]
                                df = df[1:].reset_index(drop=True)
                                table_data.append(df)

                    slide_data.append({
                        "slide_number": slide_num,
                        "content": slide_text.strip() if slide_text else None,
                        "tables": table_data or None
                    })

                pptx_data[file.replace(".pptx", "")] = slide_data
            except Exception as e:
                print(f"⚠️ Error loading PPTX {file}: {e}")

        return pptx_data if pptx_data else {}

    def ingest_all(self):
        """Load all file types and return structured data."""
        ingested_data = {
            "json": self.load_json_files(),
            "csv": self.load_csv_files(),
            "pdf": self.load_pdf_files(),
            "pptx": self.load_pptx_files()
        }
        return ingested_data


# Example usage
if __name__ == "__main__":
    ingestion = FileIngestion(verbose=True)
    data = ingestion.ingest_all()

    if ingestion.verbose:
        print("\n📌 Final Ingested Data:")

        if data["json"]:
            print("\n📂 JSON Data:")
            print(json.dumps(data["json"], indent=2, default=str))

        if data["csv"]:
            print("\n📂 CSV Data:")
            for key, df in data["csv"].items():
                print(f"📄 {key}:\n", df.head(), "\n")

        if data["pdf"]:
            print("\n📂 PDF Data:")
            for key, content in data["pdf"].items():
                print(f"📄 {key}:")
                if content["text"]:
                    print("📜 Extracted Text:", content["text"][:1])
                if content["tables"]:
                    for df in content["tables"]:
                        print("📊 Extracted Table:\n", df.head(), "\n")

        if data["pptx"]:
            print("\n📂 PPTX Data:")
            for key, slides in data["pptx"].items():
                print(f"📊 {key}:")
                for slide in slides:
                    print(f"📄 Slide {slide['slide_number']}:\n{slide['content']}\n")
                    if slide["tables"]:
                        for df in slide["tables"]:
                            print("📊 Extracted Table:\n", df.head(), "\n")
